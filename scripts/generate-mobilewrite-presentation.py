#!/usr/bin/env python3
"""Generate MobileWrite framework overview presentation for QA knowledge sharing."""

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "MobileWrite-Framework-Overview.pptx"

BLUE = RGBColor(0x00, 0x5B, 0xB5)
DARK = RGBColor(0x1A, 0x2B, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFC)
GRAY = RGBColor(0x5A, 0x6A, 0x7A)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str, subtitle=None) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(8.8), Inches(0.5))
        box.text_frame.text = subtitle
        box.text_frame.paragraphs[0].font.size = Pt(14)
        box.text_frame.paragraphs[0].font.color.rgb = GRAY


def add_bullets(slide, items: list[str], top=1.6, left=0.7, width=8.5, height=5.0) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def add_flow_diagram(slide, steps: list[str], top=2.0) -> None:
    y = top
    for i, step in enumerate(steps):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(2.5),
            Inches(y),
            Inches(5.0),
            Inches(0.55),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BLUE if i % 2 == 0 else WHITE
        shape.line.color.rgb = BLUE
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.text = step
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        y += 0.7
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, Inches(4.85), Inches(y - 0.12), Inches(0.3), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE
            arrow.line.fill.background()
            y += 0.2


def add_hierarchy(slide, lines: list[str], top=1.8) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(top), Inches(7.6), Inches(4.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BLUE
    box.line.color.rgb = BLUE
    tf = box.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Courier New"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.space_after = Pt(2)


def add_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def slide_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(2.2), Inches(10), Inches(2.8)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = BLUE
    banner.line.fill.background()
    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(1.2))
    title.text_frame.text = "MobileWrite"
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(0.8))
    sub.text_frame.text = "Mobile Automation Framework Overview"
    sub.text_frame.paragraphs[0].font.size = Pt(22)
    sub.text_frame.paragraphs[0].font.color.rgb = WHITE
    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(8.0), Inches(1.0))
    meta.text_frame.text = f"WurthLAC QE Team  |  {date.today().strftime('%B %Y')}"
    meta.text_frame.paragraphs[0].font.size = Pt(16)
    meta.text_frame.paragraphs[0].font.color.rgb = GRAY
    add_notes(
        slide,
        "Welcome the team and set context: MobileWrite is our internal automation suite built on "
        "Mobilewright for WurthLAC real-device testing (iOS and Android). Mention this session "
        "covers how the repo is organized, how to run tests, and best practices for contributors.",
    )


def slide_agenda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Agenda")
    add_bullets(
        slide,
        [
            "Introduction & purpose",
            "Why MobileWrite",
            "Key features & architecture",
            "Project structure & setup",
            "Writing & running tests",
            "Reporting & best practices",
            "Q&A",
        ],
        top=1.5,
    )
    add_notes(
        slide,
        "Walk through the agenda in under a minute. Emphasize this is a practical overview of the "
        "lac-mobile-qe repository, not a generic mobile testing lecture.",
    )


def slide_what(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "What is MobileWrite?")
    add_bullets(
        slide,
        [
            "QA automation for WurthLAC mobile app",
            "Built on Mobilewright test runner",
            "Runs on real iOS & Android devices",
            "Covers login, navigation, deals, categories",
        ],
        top=1.5,
        width=4.2,
    )
    add_flow_diagram(
        slide,
        [
            "Download EAS build",
            "Install app + agent",
            "Launch app on device",
            "Execute tests",
            "Generate report",
        ],
        top=1.7,
    )
    add_notes(
        slide,
        "MobileWrite (lac-mobile-qe) automates the WurthLAC app using Mobilewright's screen API. "
        "It solves manual regression on real devices by standardizing build install, permissions, "
        "login, and navigation flows. The workflow: fetch QA build from Expo EAS, install on a "
        "registered device, run specs from tests/, and review HTML reports on failure.",
    )


def slide_why(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Why MobileWrite?")
    icons = [
        ("⚡", "Faster test authoring"),
        ("♻️", "Reusable helpers"),
        ("🧩", "Consistent structure"),
        ("📱", "Real-device confidence"),
        ("🔧", "Easy maintenance"),
        ("📈", "Scales with new specs"),
    ]
    x_positions = [0.8, 3.5, 6.2, 0.8, 3.5, 6.2]
    y_positions = [1.8, 1.8, 1.8, 3.8, 3.8, 3.8]
    for (icon, label), x, y in zip(icons, x_positions, y_positions):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.4), Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BLUE
        card.line.color.rgb = BLUE
        tf = card.text_frame
        tf.text = f"{icon}\n{label}"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].alignment = PP_ALIGN.CENTER
            tf.paragraphs[1].font.size = Pt(12)
    add_notes(
        slide,
        "Highlight practical wins: helpers/auth.mts and helpers/navigation.mts reduce duplication; "
        "config/ centralizes bundle IDs and device IDs; one command runs Android or iOS projects. "
        "Real-device testing catches platform-specific UI issues that emulators miss.",
    )


def slide_features(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Key Features (from repo)")
    add_bullets(
        slide,
        [
            "Reusable helpers (auth, navigation, permissions)",
            "Built-in fixtures: screen, device, bundleId",
            "Environment config via .env + config/",
            "iOS & Android projects in one suite",
            "EAS build download & auto-install",
            "Retries, timeouts, HTML reporting",
            "View tree capture on failure",
        ],
        top=1.5,
    )
    add_notes(
        slide,
        "Only list what exists: helpers/ modules, Mobilewright fixtures in test files, "
        "mobilewright.config.mjs for projects/timeouts/retries, scripts/download-eas-build.mjs "
        "with --latest for Android, viewTree: on-failure in config, and npm run test:report "
        "for HTML output. No separate page-object layer—helpers serve that role.",
    )


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Framework Architecture")
    add_flow_diagram(
        slide,
        [
            "Test specs (tests/*.test.mts)",
            "Helper modules (helpers/)",
            "Config & scripts (config/, scripts/)",
            "Mobilewright runner",
            "Real device + WurthLAC app",
        ],
        top=1.6,
    )
    add_notes(
        slide,
        "Tests import helpers and call screen APIs (getByTestId, getByText, tap, fill). "
        "Config resolves bundle ID from APP_ENV. run-tests.sh sets device ID and invokes "
        "mobilewright test --project=ios|android. The Mobilewright agent on device bridges "
        "automation to the native app UI.",
    )


def slide_structure(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Project Structure")
    add_hierarchy(
        slide,
        [
            "lac-mobile-qe/",
            "├── tests/           → test specifications",
            "├── helpers/         → shared auth, navigation, permissions",
            "├── config/          → env, bundle IDs, EAS defaults",
            "├── scripts/         → download, install, run, verify",
            "├── builds/          → .ipa / .apk artifacts",
            "├── mobilewright.config.mjs",
            "└── .github/workflows/ → CI on real devices",
        ],
    )
    add_notes(
        slide,
        "Point new contributors to tests/ for specs, helpers/ for reusable logic, and config/ "
        "for environment mapping (dev/qa/production bundle IDs). builds/ is gitignored except "
        "README. scripts/ wraps common operations so QA doesn't memorize long commands.",
    )


def slide_install(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Installation & Setup")
    add_bullets(
        slide,
        [
            "Clone lac-mobile-qe repository",
            "Install dependencies",
            "Copy .env.example → .env",
            "Set device IDs & test credentials",
            "Download & install QA build",
            "Run npm run setup to verify",
        ],
        top=1.5,
    )
    add_notes(
        slide,
        "Commands for speaker: npm install; cp .env.example .env; npm run devices to find "
        "MOBILEWRIGHT_DEVICE_ID_IOS/ANDROID; npm run download:build:android && npm run "
        "install:app:android && npm run install:agent:android (one-time agent). iOS needs "
        "MOBILEWRIGHT_IOS_PROVISIONING_PROFILE for agent install. npm run setup prints a checklist.",
    )


def slide_writing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Writing Tests")
    add_bullets(
        slide,
        [
            "Import test & expect from @mobilewright/test",
            "Use screen fixture for UI actions",
            "Reuse helpers for login & navigation",
            "Assert with expect().toBeVisible()",
            "Group flows with test.step()",
        ],
        top=1.5,
        width=4.5,
    )
    code_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.8), Inches(4.2), Inches(3.8)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = DARK
    code_box.line.fill.background()
    tf = code_box.text_frame
    tf.text = (
        'test("login", async ({ screen }) => {\n'
        "  await login(screen, email, pwd);\n"
        '  await expect(screen\n'
        '    .getByText("Top Categories"))\n'
        "    .toBeVisible();\n"
        "});"
    )
    for p in tf.paragraphs:
        p.font.name = "Courier New"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
    add_notes(
        slide,
        "Show user-journey.test.mts as the advanced example: single session with test.step blocks. "
        "Prefer getByTestId, getByLabel, getByText for cross-platform stability. Keep assertions "
        "explicit—no fixed sleeps; use waitFor and expect timeouts from config.",
    )


def slide_execution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Execution Flow")
    add_flow_diagram(
        slide,
        [
            "npm run test:android",
            "run-tests.sh sets device ID",
            "Mobilewright launches app",
            "Test uses screen fixture",
            "Helper performs actions",
            "Assertions validate UI",
            "Report generated on completion",
        ],
        top=1.45,
    )
    add_notes(
        slide,
        "Android runs auto-download latest EAS build before tests (run-tests.sh). Config installs "
        "app from MOBILEWRIGHT_INSTALL_APPS_ANDROID. workers: 1 ensures one device session. "
        "Serial describe in user-journey keeps app open across steps—mimics real user flow.",
    )


def slide_reporting(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Reporting")
    items = [
        ("📊", "List reporter (default)"),
        ("🌐", "HTML report (test:report)"),
        ("🌳", "View tree on failure"),
        ("📁", "test-results/ artifacts"),
        ("☁️", "CI uploads mobilewright-report/"),
    ]
    for i, (icon, label) in enumerate(items):
        x = 0.8 + (i % 3) * 3.0
        y = 1.8 + (i // 3) * 2.0
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.6), Inches(1.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BLUE
        card.line.color.rgb = BLUE
        tf = card.text_frame
        tf.text = f"{icon}  {label}"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = DARK
    placeholder = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.8), Inches(8.0), Inches(1.8)
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    placeholder.line.color.rgb = GRAY
    placeholder.text_frame.text = "[ Screenshot placeholder: HTML report / failure view tree ]"
    placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    placeholder.text_frame.paragraphs[0].font.size = Pt(14)
    placeholder.text_frame.paragraphs[0].font.color.rgb = GRAY
    add_notes(
        slide,
        "Run npm run test:report then npm run report to open HTML. viewTree: on-failure in "
        "mobilewright.config.mjs dumps UI hierarchy for debugging. GitHub Actions workflow "
        "uploads mobilewright-report/ artifact after each run.",
    )


def slide_best_practices(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Best Practices")
    add_bullets(
        slide,
        [
            "Reuse helpers—don't duplicate login/nav logic",
            "Prefer stable locators (testID, label, text)",
            "Use test.step() for readable journeys",
            "Keep guest vs auth state in mind",
            "Run platform-specific: test:ios / test:android",
        ],
        top=1.5,
    )
    add_notes(
        slide,
        "Tips from project experience: use ensureLoggedIn/prepareApp from helpers/app-state.mts; "
        "avoid assuming guest state when device session persists; use serial mode for E2E flows; "
        "extend navigation.mts tab fallbacks instead of adding one-off taps in every test.",
    )


def slide_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "Project Summary")
    pillars = [
        ("Maintainable", "Centralized helpers & config"),
        ("Scalable", "Add specs under tests/"),
        ("Fast to extend", "Reuse existing patterns"),
        ("Real-device ready", "iOS + Android projects"),
    ]
    for i, (title, desc) in enumerate(pillars):
        x = 0.7 + (i % 2) * 4.6
        y = 1.7 + (i // 2) * 2.2
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = BLUE if i % 2 == 0 else LIGHT_BLUE
        card.line.fill.background()
        tf = card.text_frame
        tf.text = f"{title}\n{desc}"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE if i % 2 == 0 else DARK
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].font.size = Pt(13)
            tf.paragraphs[1].font.color.rgb = WHITE if i % 2 == 0 else GRAY
    add_notes(
        slide,
        "Close with the value proposition: MobileWrite gives WurthLAC QA a dedicated, version-controlled "
        "automation repo with real-device coverage, EAS integration, and a clear path for new contributors.",
    )


def slide_qa(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BLUE)
    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(2.5), Inches(1.5), Inches(5.0), Inches(3.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.text = "Questions?"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    sub = slide.shapes.add_textbox(Inches(2.0), Inches(5.2), Inches(6.0), Inches(0.6))
    sub.text_frame.text = "Thank you — WurthLAC QE Team"
    sub.text_frame.paragraphs[0].font.size = Pt(18)
    sub.text_frame.paragraphs[0].font.color.rgb = WHITE
    sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_notes(
        slide,
        "Invite questions on device setup, writing new specs, CI workflow, or extending the user "
        "journey test. Share repo link: github.com/nimeshavillvay/WLAC-WBS-App-Automation",
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_agenda(prs)
    slide_what(prs)
    slide_why(prs)
    slide_features(prs)
    slide_architecture(prs)
    slide_structure(prs)
    slide_install(prs)
    slide_writing(prs)
    slide_execution(prs)
    slide_reporting(prs)
    slide_best_practices(prs)
    slide_summary(prs)
    slide_qa(prs)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()
